"""Generate report (HTML) and presentation (.pptx) artifacts."""
import re
import secrets

from . import db, deps
from .config import ARTIFACT_DIR

REPORT_TEMPLATE = """<!doctype html>
<html><head><meta charset="utf-8"><title>{title}</title>
<style>
  body {{ font-family: Georgia, serif; max-width: 860px; margin: 3rem auto; padding: 0 1.5rem; color: #222; line-height: 1.55; }}
  h1 {{ border-bottom: 3px solid #b5542c; padding-bottom: .4rem; }}
  h2 {{ color: #b5542c; margin-top: 2rem; }}
  table {{ border-collapse: collapse; width: 100%; margin: 1rem 0; }}
  th, td {{ border: 1px solid #ccc; padding: .45rem .7rem; text-align: left; font-size: .95rem; }}
  th {{ background: #f4f1ea; }}
  .meta {{ color: #777; font-size: .85rem; }}
</style></head>
<body>
<h1>{title}</h1>
<p class="meta">Generated {created} by {author}</p>
{body}
</body></html>
"""


def _slug(title: str) -> str:
    s = re.sub(r"[^A-Za-z0-9]+", "-", title).strip("-").lower()[:40] or "artifact"
    return f"{s}-{secrets.token_hex(3)}"


def create_report(title: str, body_html: str, author: str = "") -> dict:
    filename = _slug(title) + ".html"
    html = REPORT_TEMPLATE.format(title=title, body=body_html, author=author or "customer-success bot",
                                  created=db.now())
    (ARTIFACT_DIR / filename).write_text(html)
    return db.add_artifact("report", title, filename, author)


def create_presentation(title: str, slides: list[dict], author: str = "") -> dict:
    """slides: [{"title": str, "bullets": [str, ...]}, ...]"""
    deps.require("decks")  # HTML reports work without python-pptx; decks don't
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()

    # title slide
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = f"Prepared by {author or 'customer-success bot'}"

    body_layout = prs.slide_layouts[1]
    for spec in slides:
        slide = prs.slides.add_slide(body_layout)
        slide.shapes.title.text = str(spec.get("title", ""))
        bullets = spec.get("bullets") or []
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, bullet in enumerate(bullets):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = str(bullet)
            para.font.size = Pt(20)
        notes = spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)

    filename = _slug(title) + ".pptx"
    prs.save(ARTIFACT_DIR / filename)
    return db.add_artifact("presentation", title, filename, author)
